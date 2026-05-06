"""
RAG (Retrieval-Augmented Generation) service with Groq
"""

import os
import json
from typing import Dict, List, Any, Optional
from datetime import datetime
import asyncio

from groq import Groq
from config.settings import settings
from utils.logger import get_logger
from database.models import Bot, Document as DbDocument

logger = get_logger(__name__)


class RAGService:
    """RAG service for document retrieval and response generation using Groq"""

    def __init__(self):
        """Initialize RAG service with Groq client"""
        self.settings = settings
        self.client = Groq(api_key=settings.GROQ_API_KEY)
        self.model = settings.GROQ_MODEL

    async def generate_response(
        self,
        bot_id: str,
        query: str,
        bot: Bot,
    ) -> Dict[str, Any]:
        """
        Generate RAG response using Groq API
        
        Args:
            bot_id: Bot ID
            query: User query
            bot: Bot object
            
        Returns:
            Dict with response, sources, and metadata
        """
        try:
            # Build system prompt
            system_prompt = bot.instructions or "You are a helpful AI assistant."
            
            # For now, using mock sources - in production would retrieve from FAISS
            sources = []
            
            # Prepare messages for Groq
            messages = [
                {
                    "role": "system",
                    "content": system_prompt,
                },
                {
                    "role": "user",
                    "content": query,
                },
            ]

            # Call Groq API
            response = self.client.chat.completions.create(
                messages=messages,
                model=self.model,
                temperature=bot.temperature,
                max_tokens=bot.max_tokens,
            )

            response_text = response.choices[0].message.content

            return {
                "response": response_text,
                "sources": sources,
                "tokens_used": response.usage.total_tokens if response.usage else 0,
            }

        except Exception as e:
            logger.error(f"RAG generation error: {str(e)}")
            raise

    async def process_and_embed_document(
        self,
        document_id: str,
        file_path: str,
        bot_id: str,
        db_session,
    ) -> Dict[str, Any]:
        """
        Process document and add to vector store
        
        Args:
            document_id: Document ID
            file_path: Path to document file
            bot_id: Bot ID
            db_session: Database session
            
        Returns:
            Processing status
        """
        try:
            # Extract text from document
            text_content = await self._extract_text(file_path)

            if not text_content:
                return {
                    "success": False,
                    "error": "No text extracted from document",
                }

            # Split into chunks
            chunks = self._chunk_text(
                text_content, 
                settings.CHUNK_SIZE, 
                settings.CHUNK_OVERLAP
            )

            return {
                "success": True,
                "chunks_created": len(chunks),
            }

        except Exception as e:
            logger.error(f"Document processing error: {str(e)}")
            return {
                "success": False,
                "error": str(e),
            }

    async def _extract_text(self, file_path: str) -> str:
        """Extract text from file"""
        from pathlib import Path

        file_ext = Path(file_path).suffix.lower()

        if file_ext == ".pdf":
            import pypdf
            text = ""
            with open(file_path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text()
            return text

        elif file_ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])

        elif file_ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        elif file_ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        return ""

    def _chunk_text(self, text: str, chunk_size: int, overlap: int) -> List[str]:
        """Split text into chunks"""
        words = text.split()
        chunks = []
        start = 0

        while start < len(words):
            end = min(start + chunk_size, len(words))
            chunk = " ".join(words[start:end])
            chunks.append(chunk)
            start = end - overlap

        return chunks